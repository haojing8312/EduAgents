"""
文档生成服务
自动生成教案、学生手册、评估标准等教学资料
"""

import asyncio
import json
import logging
import os
import tempfile
from abc import ABC, abstractmethod
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, BinaryIO, Dict, List, Optional, Union

import markdown

# 文档生成依赖
from jinja2 import Environment, FileSystemLoader, select_autoescape
from markdown.extensions import codehilite, tables, toc
from pptx import Presentation
from pptx.util import Inches as PptxInches
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Inches, Pt
from weasyprint import CSS, HTML

from ..core.config import settings
from ..models.course import Assessment, Course, Lesson, Resource

logger = logging.getLogger(__name__)


class DocumentGeneratorError(Exception):
    """文档生成异常"""
    pass


class BaseDocumentGenerator(ABC):
    """文档生成器基类"""
    
    def __init__(self, template_dir: str = None):
        self.template_dir = template_dir or os.path.join(
            os.path.dirname(__file__), "..", "..", "templates", "documents"
        )
        self.jinja_env = Environment(
            loader=FileSystemLoader(self.template_dir),
            autoescape=select_autoescape(['html', 'xml'])
        )
        self._setup_jinja_filters()
    
    def _setup_jinja_filters(self):
        """设置Jinja2过滤器"""
        self.jinja_env.filters['format_duration'] = self._format_duration
        self.jinja_env.filters['format_date'] = self._format_date
        self.jinja_env.filters['markdown'] = self._markdown_filter
        self.jinja_env.filters['truncate_words'] = self._truncate_words
    
    def _format_duration(self, minutes: int) -> str:
        """格式化时长"""
        if minutes < 60:
            return f"{minutes}分钟"
        hours = minutes // 60
        remaining_minutes = minutes % 60
        if remaining_minutes == 0:
            return f"{hours}小时"
        return f"{hours}小时{remaining_minutes}分钟"
    
    def _format_date(self, date: datetime, format_str: str = "%Y年%m月%d日") -> str:
        """格式化日期"""
        return date.strftime(format_str)
    
    def _markdown_filter(self, text: str) -> str:
        """Markdown转HTML过滤器"""
        if not text:
            return ""
        return markdown.markdown(
            text, 
            extensions=['codehilite', 'tables', 'toc', 'fenced_code']
        )
    
    def _truncate_words(self, text: str, length: int = 100) -> str:
        """截断文本"""
        if not text or len(text) <= length:
            return text
        return text[:length] + "..."
    
    @abstractmethod
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成文档"""
        pass
    
    def _prepare_course_data(self, course: Course) -> Dict[str, Any]:
        """准备课程数据"""
        return {
            'course': course,
            'generated_at': datetime.now(),
            'generator_info': {
                'name': self.__class__.__name__,
                'version': '1.0.0'
            }
        }


class TeachingPlanGenerator(BaseDocumentGenerator):
    """教案生成器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成教案文档"""
        try:
            options = options or {}
            template_name = options.get('template', 'teaching_plan.html')
            
            # 准备数据
            data = self._prepare_course_data(course)
            data.update({
                'include_objectives': options.get('include_objectives', True),
                'include_activities': options.get('include_activities', True),
                'include_assessments': options.get('include_assessments', True),
                'include_resources': options.get('include_resources', True),
                'detailed_lessons': options.get('detailed_lessons', True)
            })
            
            # 渲染模板
            template = self.jinja_env.get_template(template_name)
            html_content = template.render(**data)
            
            # 生成PDF
            css_path = os.path.join(self.template_dir, 'styles', 'teaching_plan.css')
            css = CSS(css_path) if os.path.exists(css_path) else None
            
            pdf_buffer = BytesIO()
            HTML(string=html_content).write_pdf(pdf_buffer, stylesheets=[css] if css else None)
            
            return pdf_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"生成教案失败: {str(e)}")
            raise DocumentGeneratorError(f"生成教案失败: {str(e)}")


class StudentHandbookGenerator(BaseDocumentGenerator):
    """学生手册生成器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成学生手册"""
        try:
            options = options or {}
            template_name = options.get('template', 'student_handbook.html')
            
            # 准备数据
            data = self._prepare_course_data(course)
            data.update({
                'include_schedule': options.get('include_schedule', True),
                'include_assignments': options.get('include_assignments', True),
                'include_rubrics': options.get('include_rubrics', True),
                'include_resources': options.get('include_resources', True),
                'student_friendly': True
            })
            
            # 渲染模板
            template = self.jinja_env.get_template(template_name)
            html_content = template.render(**data)
            
            # 生成PDF
            css_path = os.path.join(self.template_dir, 'styles', 'student_handbook.css')
            css = CSS(css_path) if os.path.exists(css_path) else None
            
            pdf_buffer = BytesIO()
            HTML(string=html_content).write_pdf(pdf_buffer, stylesheets=[css] if css else None)
            
            return pdf_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"生成学生手册失败: {str(e)}")
            raise DocumentGeneratorError(f"生成学生手册失败: {str(e)}")


class AssessmentRubricGenerator(BaseDocumentGenerator):
    """评估量规生成器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成评估量规"""
        try:
            options = options or {}
            template_name = options.get('template', 'assessment_rubric.html')
            
            # 准备数据
            data = self._prepare_course_data(course)
            data.update({
                'rubric_style': options.get('rubric_style', 'detailed'),
                'include_examples': options.get('include_examples', True),
                'scoring_method': options.get('scoring_method', 'points')
            })
            
            # 渲染模板
            template = self.jinja_env.get_template(template_name)
            html_content = template.render(**data)
            
            # 生成PDF
            css_path = os.path.join(self.template_dir, 'styles', 'assessment_rubric.css')
            css = CSS(css_path) if os.path.exists(css_path) else None
            
            pdf_buffer = BytesIO()
            HTML(string=html_content).write_pdf(pdf_buffer, stylesheets=[css] if css else None)
            
            return pdf_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"生成评估量规失败: {str(e)}")
            raise DocumentGeneratorError(f"生成评估量规失败: {str(e)}")


class WordDocumentGenerator(BaseDocumentGenerator):
    """Word文档生成器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成Word文档"""
        try:
            options = options or {}
            doc_type = options.get('type', 'complete')  # complete, teaching_plan, handbook
            
            # 创建文档
            document = Document()
            
            # 设置文档样式
            self._setup_document_styles(document)
            
            if doc_type == 'complete':
                await self._generate_complete_document(document, course, options)
            elif doc_type == 'teaching_plan':
                await self._generate_teaching_plan_document(document, course, options)
            elif doc_type == 'handbook':
                await self._generate_handbook_document(document, course, options)
            
            # 保存到内存
            doc_buffer = BytesIO()
            document.save(doc_buffer)
            
            return doc_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"生成Word文档失败: {str(e)}")
            raise DocumentGeneratorError(f"生成Word文档失败: {str(e)}")
    
    def _setup_document_styles(self, document: Document):
        """设置文档样式"""
        # 设置标题样式
        title_style = document.styles['Title']
        title_style.font.size = Pt(18)
        title_style.font.bold = True
        
        # 设置正文样式
        normal_style = document.styles['Normal']
        normal_style.font.size = Pt(12)
        normal_style.paragraph_format.line_spacing = 1.15
    
    async def _generate_complete_document(self, document: Document, course: Course, options: Dict[str, Any]):
        """生成完整文档"""
        # 添加标题
        title = document.add_heading(course.title, 0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        if course.subtitle:
            subtitle = document.add_paragraph(course.subtitle)
            subtitle.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # 添加课程概述
        document.add_heading('课程概述', level=1)
        if course.description:
            document.add_paragraph(course.description)
        
        # 添加学习目标
        if course.learning_objectives:
            document.add_heading('学习目标', level=1)
            for i, objective in enumerate(course.learning_objectives, 1):
                document.add_paragraph(f"{i}. {objective}", style='List Number')
        
        # 添加课程结构
        if course.lessons:
            document.add_heading('课程结构', level=1)
            for lesson in course.lessons:
                document.add_heading(lesson.title, level=2)
                if lesson.description:
                    document.add_paragraph(lesson.description)
                
                # 添加学习活动
                if lesson.activities:
                    document.add_heading('学习活动', level=3)
                    for activity in lesson.activities:
                        document.add_paragraph(f"• {activity}", style='List Bullet')
        
        # 添加评估方式
        if course.assessments:
            document.add_heading('评估方式', level=1)
            for assessment in course.assessments:
                document.add_heading(assessment.title, level=2)
                if assessment.description:
                    document.add_paragraph(assessment.description)
    
    async def _generate_teaching_plan_document(self, document: Document, course: Course, options: Dict[str, Any]):
        """生成教案文档"""
        # 实现教案生成逻辑
        pass
    
    async def _generate_handbook_document(self, document: Document, course: Course, options: Dict[str, Any]):
        """生成手册文档"""
        # 实现手册生成逻辑
        pass


class PowerPointGenerator(BaseDocumentGenerator):
    """PowerPoint演示文稿生成器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """生成PowerPoint演示文稿"""
        try:
            options = options or {}
            
            # 创建演示文稿
            presentation = Presentation()
            
            # 设置幻灯片尺寸（16:9）
            presentation.slide_width = PptxInches(13.33)
            presentation.slide_height = PptxInches(7.5)
            
            # 生成封面页
            self._add_title_slide(presentation, course)
            
            # 生成课程概述页
            self._add_overview_slide(presentation, course)
            
            # 生成学习目标页
            if course.learning_objectives:
                self._add_objectives_slide(presentation, course)
            
            # 生成课程结构页
            if course.lessons:
                self._add_structure_slide(presentation, course)
                
                # 为每个课时生成幻灯片
                for lesson in course.lessons:
                    self._add_lesson_slide(presentation, lesson)
            
            # 生成评估页
            if course.assessments:
                self._add_assessment_slide(presentation, course)
            
            # 生成结束页
            self._add_conclusion_slide(presentation, course)
            
            # 保存到内存
            ppt_buffer = BytesIO()
            presentation.save(ppt_buffer)
            
            return ppt_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"生成PowerPoint失败: {str(e)}")
            raise DocumentGeneratorError(f"生成PowerPoint失败: {str(e)}")
    
    def _add_title_slide(self, presentation: Presentation, course: Course):
        """添加标题页"""
        slide_layout = presentation.slide_layouts[0]  # 标题页布局
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = course.title
        subtitle.text = course.subtitle or course.description[:100] + "..." if course.description else ""
    
    def _add_overview_slide(self, presentation: Presentation, course: Course):
        """添加概述页"""
        slide_layout = presentation.slide_layouts[1]  # 标题和内容布局
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "课程概述"
        
        # 构建概述内容
        overview_text = []
        if course.description:
            overview_text.append(f"课程描述：{course.description}")
        
        overview_text.append(f"适用学段：{course.education_level}")
        overview_text.append(f"主要学科：{course.subject}")
        overview_text.append(f"课程周数：{course.duration_weeks}周")
        overview_text.append(f"总学时：{course.duration_hours}小时")
        
        content.text = "\n".join(overview_text)
    
    def _add_objectives_slide(self, presentation: Presentation, course: Course):
        """添加学习目标页"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "学习目标"
        
        objectives_text = []
        for i, objective in enumerate(course.learning_objectives, 1):
            objectives_text.append(f"{i}. {objective}")
        
        content.text = "\n".join(objectives_text)
    
    def _add_structure_slide(self, presentation: Presentation, course: Course):
        """添加课程结构页"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "课程结构"
        
        structure_text = []
        for i, lesson in enumerate(course.lessons, 1):
            structure_text.append(f"{i}. {lesson.title}")
        
        content.text = "\n".join(structure_text)
    
    def _add_lesson_slide(self, presentation: Presentation, lesson: Lesson):
        """添加课时页"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = lesson.title
        
        lesson_text = []
        if lesson.description:
            lesson_text.append(f"描述：{lesson.description}")
        
        lesson_text.append(f"时长：{lesson.duration_minutes}分钟")
        
        if lesson.objectives:
            lesson_text.append("学习目标：")
            for objective in lesson.objectives:
                lesson_text.append(f"• {objective}")
        
        content.text = "\n".join(lesson_text)
    
    def _add_assessment_slide(self, presentation: Presentation, course: Course):
        """添加评估页"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "评估方式"
        
        assessment_text = []
        for assessment in course.assessments:
            assessment_text.append(f"• {assessment.title}")
            if assessment.description:
                assessment_text.append(f"  {assessment.description}")
        
        content.text = "\n".join(assessment_text)
    
    def _add_conclusion_slide(self, presentation: Presentation, course: Course):
        """添加结束页"""
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "谢谢"
        subtitle.text = f"{course.title} - 课程介绍完毕"


class JSONExporter(BaseDocumentGenerator):
    """JSON导出器"""
    
    async def generate(self, course: Course, options: Dict[str, Any] = None) -> bytes:
        """导出为JSON格式"""
        try:
            options = options or {}
            include_metadata = options.get('include_metadata', True)
            indent = options.get('indent', 2)
            
            # 构建导出数据
            export_data = {
                'course': self._serialize_course(course),
                'export_info': {
                    'generated_at': datetime.now().isoformat(),
                    'generator': 'JSONExporter',
                    'version': '1.0.0',
                    'options': options
                } if include_metadata else None
            }
            
            # 转换为JSON
            json_str = json.dumps(export_data, ensure_ascii=False, indent=indent)
            
            return json_str.encode('utf-8')
            
        except Exception as e:
            logger.error(f"导出JSON失败: {str(e)}")
            raise DocumentGeneratorError(f"导出JSON失败: {str(e)}")
    
    def _serialize_course(self, course: Course) -> Dict[str, Any]:
        """序列化课程对象"""
        data = course.to_dict()
        
        # 序列化关联对象
        if course.lessons:
            data['lessons'] = [lesson.to_dict() for lesson in course.lessons]
        
        if course.assessments:
            data['assessments'] = [assessment.to_dict() for assessment in course.assessments]
        
        if course.resources:
            data['resources'] = [resource.to_dict() for resource in course.resources]
        
        return data


class DocumentGeneratorService:
    """文档生成服务"""
    
    def __init__(self):
        self.generators = {
            'pdf_teaching_plan': TeachingPlanGenerator(),
            'pdf_handbook': StudentHandbookGenerator(),
            'pdf_rubric': AssessmentRubricGenerator(),
            'docx': WordDocumentGenerator(),
            'pptx': PowerPointGenerator(),
            'json': JSONExporter()
        }
    
    async def generate_document(
        self,
        course: Course,
        format_type: str,
        options: Dict[str, Any] = None
    ) -> bytes:
        """生成文档"""
        if format_type not in self.generators:
            raise DocumentGeneratorError(f"不支持的格式: {format_type}")
        
        generator = self.generators[format_type]
        return await generator.generate(course, options)
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的格式列表"""
        return list(self.generators.keys())
    
    async def generate_complete_package(
        self,
        course: Course,
        formats: List[str] = None,
        options: Dict[str, Any] = None
    ) -> Dict[str, bytes]:
        """生成完整的教学资料包"""
        formats = formats or ['pdf_teaching_plan', 'pdf_handbook', 'docx', 'json']
        options = options or {}
        
        results = {}
        
        for format_type in formats:
            if format_type in self.generators:
                try:
                    content = await self.generate_document(course, format_type, options.get(format_type, {}))
                    results[format_type] = content
                except Exception as e:
                    logger.error(f"生成{format_type}格式失败: {str(e)}")
                    # 继续生成其他格式，不中断整个流程
        
        return results


# 全局服务实例
document_generator_service = DocumentGeneratorService()